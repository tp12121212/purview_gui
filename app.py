import os
import csv
import json
import tempfile
import sys
import hashlib
import sqlite3
import tarfile
import zipfile
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import re
import html as html_lib
from email import policy
from email.parser import BytesParser
import xml.etree.ElementTree as ET
import base64

import pytesseract
from PIL import Image, ImageFile, ImageOps, ImageFilter, ExifTags, ImageSequence
try:
    import cv2
    import numpy as np
except Exception:  # Optional for OCR preprocessing
    cv2 = None
    np = None
try:
    import easyocr
except Exception:  # Optional OCR backend
    easyocr = None
import PyPDF2
import docx
import openpyxl
from pdf2image import convert_from_path
import extract_msg
import warnings
import rarfile
import py7zr
from pptx import Presentation

if sys.version_info[:2] != (3, 11):
    raise SystemExit("Python 3.11 is required. Please use a 3.11 virtual environment.")

# ======================================================
# Flask app
# ======================================================

app = Flask(__name__)
CORS(app, resources={
    r"/*": {
        "origins": [
            "http://127.0.0.1:5500",
            "http://localhost:5500",
            "null",
            "file://",
        ]
    }
})

# Allow large images while keeping a sane cap to avoid DoS warnings.
Image.MAX_IMAGE_PIXELS = 200_000_000
warnings.filterwarnings(
    "ignore",
    category=Image.DecompressionBombWarning,
)
ImageFile.LOAD_TRUNCATED_IMAGES = True

# ======================================================
# Configuration
# ======================================================

DEFAULT_LEXICON_PATH = "lexicon_latest.csv"
REGEX_PATTERNS_PATH = "regex_patterns.py"
PROGRESS_DIR = Path(tempfile.gettempdir()) / "purview_scan_progress"
RULEPACK_CACHE_PATH = Path("rulepack_cache.json")
MAX_SIT_FILES = 50
DB_PATH = Path(os.environ.get("PURVIEW_DB_PATH", "scan_results.db"))

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".jpg", ".jpeg", ".png", ".tif", ".tiff",
    ".eml", ".msg",
    ".zip", ".7z", ".rar", ".tar", ".tgz", ".tar.gz"
}

app.config["MAX_CONTENT_LENGTH"] = 500 * 1024 * 1024  # 500MB
MAX_NESTED_DEPTH = 2

# ======================================================
# Database
# ======================================================

def get_db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS scans (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            created_at TEXT NOT NULL,
            scenario_mode TEXT NOT NULL,
            status TEXT NOT NULL,
            scan_params TEXT NOT NULL
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            scan_id INTEGER NOT NULL,
            source_type TEXT NOT NULL,
            source_uri TEXT NOT NULL,
            sha256 TEXT NOT NULL,
            extracted_text_ref TEXT,
            FOREIGN KEY(scan_id) REFERENCES scans(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS regex_patterns (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            pattern TEXT NOT NULL,
            version TEXT,
            source_file TEXT
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS lexicon_entries (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            keyword TEXT NOT NULL,
            type TEXT,
            source_file TEXT
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS matches (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            scan_id INTEGER NOT NULL,
            document_id INTEGER NOT NULL,
            match_type TEXT NOT NULL,
            pattern_name TEXT,
            value TEXT,
            position INTEGER,
            context TEXT,
            nearest_regex TEXT,
            distance INTEGER,
            container_path TEXT,
            inner_path TEXT,
            FOREIGN KEY(scan_id) REFERENCES scans(id),
            FOREIGN KEY(document_id) REFERENCES documents(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS sit_definitions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            description TEXT,
            publisher TEXT,
            created_from_scan_id INTEGER NOT NULL,
            rule_pack_name TEXT NOT NULL,
            version TEXT,
            FOREIGN KEY(created_from_scan_id) REFERENCES scans(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS rule_packs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            sit_definition_id INTEGER NOT NULL,
            xml_blob TEXT NOT NULL,
            generated_at TEXT NOT NULL,
            generation_params TEXT NOT NULL,
            FOREIGN KEY(sit_definition_id) REFERENCES sit_definitions(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS labels (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            label_id TEXT,
            label_name TEXT,
            source TEXT,
            applied_at TEXT,
            status TEXT,
            FOREIGN KEY(document_id) REFERENCES documents(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS scan_policies (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            config_json TEXT NOT NULL
        )
        """
    )
    conn.commit()
    conn.close()


init_db()

# ======================================================
# Lexicon loading (single-word only)
# ======================================================

def load_keywords(lexicon_path: str = DEFAULT_LEXICON_PATH, allowed_types=None):
    keywords = []

    with open(lexicon_path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        type_key = None
        if reader.fieldnames:
            for name in reader.fieldnames:
                if name and name.strip().lower() == "type":
                    type_key = name
                    break
        for row in reader:
            value = (
                row.get("keyword")
                or row.get("phrase")
                or row.get("value")
            )
            if not value:
                continue

            value = value.strip()

            # Single word, alnum only
            if re.fullmatch(r"[A-Za-z0-9]+", value):
                if allowed_types and type_key:
                    row_type = (row.get(type_key) or "").strip().lower() or "other"
                    if row_type not in allowed_types:
                        continue
                keywords.append(value.lower())

    return keywords

def load_lexicon_entries(lexicon_path: str = DEFAULT_LEXICON_PATH):
    entries = []
    with open(lexicon_path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        type_key = None
        if reader.fieldnames:
            for name in reader.fieldnames:
                if name and name.strip().lower() == "type":
                    type_key = name
                    break
        for row in reader:
            value = (
                row.get("keyword")
                or row.get("phrase")
                or row.get("value")
            )
            if not value:
                continue
            value = value.strip()
            if re.fullmatch(r"[A-Za-z0-9]+", value):
                entry_type = None
                if type_key:
                    entry_type = (row.get(type_key) or "").strip() or None
                entries.append({"keyword": value.lower(), "type": entry_type})
    return entries

# ======================================================
# Regex loading
# ======================================================

def load_regex_patterns(regex_path: str = REGEX_PATTERNS_PATH):
    path = Path(regex_path)
    if path.suffix.lower() == ".py":
        import importlib.util
        spec = importlib.util.spec_from_file_location("regex_patterns", path)
        if spec is None or spec.loader is None:
            raise ValueError(f"unable to load regex patterns from {regex_path}")
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)
        patterns = getattr(module, "patterns", None)
        if patterns is None:
            raise ValueError(f"regex patterns missing 'patterns' in {regex_path}")
        return patterns
    with open(regex_path, encoding="utf-8") as f:
        return json.load(f)

def upsert_regex_patterns(conn, patterns, source_file: str):
    cursor = conn.cursor()
    cursor.execute("DELETE FROM regex_patterns WHERE source_file = ?", (source_file,))
    for item in patterns:
        cursor.execute(
            """
            INSERT INTO regex_patterns (name, pattern, version, source_file)
            VALUES (?, ?, ?, ?)
            """,
            (
                item.get("name", "Regex"),
                item.get("pattern", ""),
                item.get("version"),
                source_file,
            ),
        )
    conn.commit()

def upsert_lexicon_entries(conn, entries, source_file: str):
    cursor = conn.cursor()
    cursor.execute("DELETE FROM lexicon_entries WHERE source_file = ?", (source_file,))
    for entry in entries:
        cursor.execute(
            """
            INSERT INTO lexicon_entries (keyword, type, source_file)
            VALUES (?, ?, ?)
            """,
            (entry["keyword"], entry.get("type"), source_file),
        )
    conn.commit()

def sha256_for_path(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()

# ======================================================
# File collection
# ======================================================

def collect_files_from_path(path: str, recursive: bool):
    collected = []
    p = Path(path)

    if p.is_file():
        if p.suffix.lower() in SUPPORTED_EXTENSIONS:
            collected.append(p)
        return collected

    if p.is_dir():
        iterator = p.rglob("*") if recursive else p.glob("*")
        for item in iterator:
            if item.is_file() and item.suffix.lower() in SUPPORTED_EXTENSIONS:
                collected.append(item)

    return collected

def parse_file_types(value: str):
    if not value:
        return None
    types = set()
    for item in value.split(","):
        ext = item.strip().lower()
        if not ext:
            continue
        if not ext.startswith("."):
            ext = "." + ext
        types.add(ext)
    return types or None

def parse_lexicon_types(value: str):
    if not value:
        return None
    types = set()
    for item in value.split(","):
        label = item.strip().lower()
        if label:
            types.add(label)
    return types or None

def allowed_extension(ext: str, allowed_exts):
    if allowed_exts is None:
        return ext in SUPPORTED_EXTENSIONS
    return ext in allowed_exts

# ======================================================
# Scan iteration
# ======================================================

def iter_scan_items(files, scan_path: str, recursive: bool, allowed_exts):
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if allowed_extension(ext, allowed_exts):
            yield ("upload", f)

    if scan_path:
        for p in collect_files_from_path(scan_path, recursive):
            if allowed_extension(p.suffix.lower(), allowed_exts):
                yield ("path", p)

# ======================================================
# Text extraction
# ======================================================

def normalize_image_for_ocr(image: Image.Image) -> Image.Image:
    try:
        image = ImageOps.exif_transpose(image)
    except Exception:
        pass
    if image.mode not in ("RGB", "L"):
        image = image.convert("RGB")
    return image

def extract_exif_text(image: Image.Image) -> str:
    try:
        exif = image.getexif()
    except Exception:
        exif = None
    if not exif:
        return ""
    parts = []
    tag_map = ExifTags.TAGS
    orientation_map = {
        1: "Top side, left",
        2: "Top side, right (Mirror horizontal)",
        3: "Bottom side, right (Rotate 180)",
        4: "Bottom side, left (Mirror vertical)",
        5: "Left side, top (Mirror horizontal and rotate 270 CW)",
        6: "Right side, top (Rotate 90 CW)",
        7: "Right side, bottom (Mirror horizontal and rotate 90 CW)",
        8: "Left side, bottom (Rotate 270 CW)",
    }
    for key, value in exif.items():
        name = tag_map.get(key, str(key))
        if name == "Orientation":
            value = orientation_map.get(value, value)
        parts.append(f"{name}: {value}")
    return " ".join(parts)

def order_points(pts):
    rect = np.zeros((4, 2), dtype="float32")
    s = pts.sum(axis=1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]
    diff = np.diff(pts, axis=1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    return rect

def four_point_transform(image, pts):
    rect = order_points(pts)
    (tl, tr, br, bl) = rect
    widthA = np.linalg.norm(br - bl)
    widthB = np.linalg.norm(tr - tl)
    maxWidth = int(max(widthA, widthB))
    heightA = np.linalg.norm(tr - br)
    heightB = np.linalg.norm(tl - bl)
    maxHeight = int(max(heightA, heightB))
    dst = np.array([
        [0, 0],
        [maxWidth - 1, 0],
        [maxWidth - 1, maxHeight - 1],
        [0, maxHeight - 1]
    ], dtype="float32")
    M = cv2.getPerspectiveTransform(rect, dst)
    warped = cv2.warpPerspective(image, M, (maxWidth, maxHeight))
    return warped

def extract_document_perspective(image: Image.Image) -> Image.Image | None:
    if cv2 is None or np is None:
        return None
    img = np.array(image.convert("RGB"))
    orig = img.copy()
    height, width = img.shape[:2]
    max_dim = max(height, width)
    ratio = 1.0
    if max_dim > 1200:
        ratio = max_dim / 1200.0
        img = cv2.resize(img, (int(width / ratio), int(height / ratio)), interpolation=cv2.INTER_AREA)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(gray, (5, 5), 0)
    edged = cv2.Canny(blur, 50, 150)
    cnts, _ = cv2.findContours(edged, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    cnts = sorted(cnts, key=cv2.contourArea, reverse=True)[:10]
    for c in cnts:
        peri = cv2.arcLength(c, True)
        approx = cv2.approxPolyDP(c, 0.02 * peri, True)
        if len(approx) == 4 and cv2.contourArea(approx) > 5000:
            pts = approx.reshape(4, 2).astype("float32") * ratio
            warped = four_point_transform(orig, pts)
            return Image.fromarray(warped)
    return None

def select_best_frame(image: Image.Image) -> Image.Image:
    if getattr(image, "format", "").upper() != "MPO":
        return image
    best = image
    best_score = -1.0
    for frame in ImageSequence.Iterator(image):
        try:
            frame = frame.copy()
        except Exception:
            continue
        if cv2 is not None and np is not None:
            arr = np.array(frame)
            if arr.ndim == 2:
                gray = arr
            else:
                gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
            score = cv2.Laplacian(gray, cv2.CV_64F).var()
        else:
            edges = frame.convert("L").filter(ImageFilter.FIND_EDGES)
            score = sum(edges.getdata()) / max(1, edges.size[0] * edges.size[1])
        if score > best_score:
            best = frame
            best_score = score
    return best

def is_macos() -> bool:
    return sys.platform == "darwin"

def ensure_vision_ocr_binary() -> Path | None:
    if not is_macos():
        return None
    bin_path = Path("scripts/vision_ocr")
    if bin_path.exists() and os.access(bin_path, os.X_OK):
        return bin_path
    swift_path = Path("scripts/vision_ocr.swift")
    if not swift_path.exists():
        return None
    try:
        result = subprocess.run(
            ["xcrun", "swiftc", "-O", str(swift_path), "-o", str(bin_path)],
            capture_output=True,
            text=True,
            check=True,
        )
        if bin_path.exists():
            return bin_path
    except Exception as e:
        print(f"Vision OCR build failed: {e}")
    return None

def vision_ocr_from_path(path: Path) -> str:
    bin_path = ensure_vision_ocr_binary()
    if not bin_path:
        raise RuntimeError("Vision OCR binary unavailable")
    result = subprocess.run(
        [str(bin_path), str(path)],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or "Vision OCR failed")
    return result.stdout.strip()

def vision_ocr_from_image(image: Image.Image) -> str:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp_name = tmp.name
        image.save(tmp_name, format="PNG")
    try:
        return vision_ocr_from_path(Path(tmp_name))
    finally:
        try:
            os.unlink(tmp_name)
        except OSError:
            pass

_easyocr_reader = None

def get_easyocr_reader():
    global _easyocr_reader
    if easyocr is None:
        return None
    if _easyocr_reader is None:
        _easyocr_reader = easyocr.Reader(["en"], gpu=False)
    return _easyocr_reader

def easyocr_from_image(image: Image.Image) -> str:
    reader = get_easyocr_reader()
    if reader is None:
        raise RuntimeError("EasyOCR unavailable")
    if np is None:
        raise RuntimeError("NumPy unavailable for EasyOCR")
    img = np.array(image)
    results = reader.readtext(img, detail=0, paragraph=True)
    if isinstance(results, (list, tuple)):
        return "\n".join(str(r) for r in results if r)
    return str(results).strip()

def apply_osd_rotation(image: Image.Image) -> Image.Image:
    try:
        probe = image
        max_dim = max(image.width, image.height)
        if max_dim > 1600:
            scale = 1600 / max_dim
            probe = image.resize(
                (int(image.width * scale), int(image.height * scale)),
                Image.LANCZOS,
            )
        osd = pytesseract.image_to_osd(probe, output_type=pytesseract.Output.DICT)
        rotate = int(osd.get("rotate", 0) or 0)
    except Exception:
        return image
    if rotate in (90, 180, 270):
        return image.rotate(-rotate, expand=True)
    return image

def ocr_image_confidence(image, lang: str = "eng", psm: int = 6):
    config = f"--oem 3 --psm {psm}"
    if isinstance(image, Image.Image):
        image = normalize_image_for_ocr(image)
    data = pytesseract.image_to_data(image, lang=lang, config=config, output_type=pytesseract.Output.DICT)
    confidences = [int(c) for c in data.get("conf", []) if c and c != "-1"]
    if not confidences:
        return None
    return sum(confidences) / len(confidences)

def score_ocr_text(text: str, conf: float = 0.0) -> float:
    if not text:
        return 0.0
    cleaned = re.sub(r"[^\w\s]", " ", text)
    tokens = re.findall(r"[A-Za-z0-9]+", cleaned)
    if not tokens:
        return 0.0
    word_tokens = [t for t in tokens if any(c.isalpha() for c in t)]
    word_count = len(word_tokens)
    avg_len = sum(len(t) for t in word_tokens) / max(1, word_count)
    alpha_ratio = len([t for t in word_tokens if t.isalpha()]) / max(1, word_count)
    vowel_ratio = len([t for t in word_tokens if re.search(r"[aeiouAEIOU]", t)]) / max(1, word_count)
    digit_ratio = len([t for t in tokens if t.isdigit()]) / max(1, len(tokens))
    weird_ratio = len(re.findall(r"[^\x20-\x7E\n\r\t]", text)) / max(1, len(text))
    score = (word_count * 2.0) + avg_len + (alpha_ratio * 5.0) + (vowel_ratio * 5.0) + (digit_ratio * 1.0)
    score -= weird_ratio * 10.0
    score += conf * 0.1
    return score

def ocr_text_from_regions(image: Image.Image, lang: str = "eng") -> str:
    if cv2 is None or np is None:
        return ""
    gray = np.array(image.convert("L"))
    height, width = gray.shape[:2]
    max_dim = max(height, width)
    if max_dim > 2400:
        scale = 2400 / max_dim
        gray = cv2.resize(gray, (int(width * scale), int(height * scale)), interpolation=cv2.INTER_AREA)
        height, width = gray.shape[:2]

    blur = cv2.GaussianBlur(gray, (3, 3), 0)
    _, bw = cv2.threshold(blur, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    if np.mean(bw) > 127:
        bw = cv2.bitwise_not(bw)

    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (25, 3))
    connected = cv2.morphologyEx(bw, cv2.MORPH_CLOSE, kernel, iterations=1)
    contours, _ = cv2.findContours(connected, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    boxes = []
    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        if h < 12 or w < 60:
            continue
        if w * h < 800:
            continue
        boxes.append((y, x, w, h))

    if not boxes:
        return ""
    boxes.sort()
    lines = []
    for y, x, w, h in boxes:
        pad = 4
        x0 = max(0, x - pad)
        y0 = max(0, y - pad)
        x1 = min(width, x + w + pad)
        y1 = min(height, y + h + pad)
        crop = gray[y0:y1, x0:x1]
        crop_img = Image.fromarray(crop)
        psm = 7 if h < 40 else 6
        config = f"--oem 3 --psm {psm} --dpi 300"
        try:
            line = pytesseract.image_to_string(crop_img, lang=lang, config=config)
        except pytesseract.TesseractError:
            line = ""
        line = line.strip()
        if line:
            lines.append(line)

    return "\n".join(lines)

def ocr_image_to_text(image, lang: str = "eng", source_path: Path | None = None):
    backend = os.environ.get("OCR_BACKEND", "auto").strip().lower()
    if backend in {"vision", "auto"}:
        bin_path = ensure_vision_ocr_binary() if backend == "auto" else ensure_vision_ocr_binary()
        if bin_path:
            try:
                if source_path:
                    return vision_ocr_from_path(source_path)
                if isinstance(image, Image.Image):
                    return vision_ocr_from_image(image)
            except Exception as e:
                print(f"Vision OCR fallback to Tesseract: {e}")

    if backend in {"easyocr", "auto"}:
        try:
            if isinstance(image, Image.Image):
                return easyocr_from_image(image)
        except Exception as e:
            if backend == "easyocr":
                print(f"EasyOCR failed: {e}")
            else:
                print(f"EasyOCR fallback to Tesseract: {e}")

    if isinstance(image, Image.Image):
        image = normalize_image_for_ocr(image)
        if os.environ.get("OCR_USE_OSD", "true").strip().lower() in {"1", "true", "yes"}:
            image = apply_osd_rotation(image)

    def run_ocr(img, psm: int):
        config = f"--oem 3 --psm {psm} --dpi 300"
        try:
            text = pytesseract.image_to_string(img, lang=lang, config=config)
            conf = ocr_image_confidence(img, lang=lang, psm=psm) or 0
            word_count = len(re.findall(r"\w+", text))
            return text, conf, word_count
        except pytesseract.TesseractError:
            if isinstance(img, Image.Image):
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp_name = tmp.name
                    img.save(tmp_name, format="PNG")
                try:
                    text = pytesseract.image_to_string(tmp_name, lang=lang, config=config)
                    conf = ocr_image_confidence(tmp_name, lang=lang, psm=psm) or 0
                    word_count = len(re.findall(r"\w+", text))
                    return text, conf, word_count
                finally:
                    try:
                        os.unlink(tmp_name)
                    except OSError:
                        pass
            raise

    def ocr_from_base(base_img: Image.Image):
        gray_full = ImageOps.autocontrast(base_img.convert("L"))

        # Fast pass on downscaled image
        max_dim_full = max(gray_full.width, gray_full.height)
        fast = gray_full
        if max_dim_full > 1600:
            scale = 1600 / max_dim_full
            fast = gray_full.resize(
                (int(gray_full.width * scale), int(gray_full.height * scale)),
                Image.LANCZOS,
            )

        fast_candidates = [fast, fast.filter(ImageFilter.SHARPEN)]
        best_text = ""
        best_score = 0.0
        for img in fast_candidates:
            text, conf, _ = run_ocr(img, psm=6)
            score = score_ocr_text(text, conf=conf)
            if score > best_score:
                best_score = score
                best_text = text

        # Full-quality fallback
        candidates = [gray_full, gray_full.filter(ImageFilter.SHARPEN)]
        if max_dim_full < 2200:
            candidates.append(
                gray_full.resize((gray_full.width * 2, gray_full.height * 2), Image.LANCZOS)
            )
        if cv2 is not None and np is not None:
            try:
                cv_gray = np.array(gray_full)
                blur = cv2.bilateralFilter(cv_gray, 9, 75, 75)
                adap = cv2.adaptiveThreshold(
                    blur, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 31, 5
                )
                kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (2, 2))
                cleaned = cv2.morphologyEx(adap, cv2.MORPH_CLOSE, kernel)
                candidates.append(Image.fromarray(cleaned))
            except Exception:
                pass
        else:
            bw = gray_full.point(lambda x: 0 if x < 160 else 255, mode="1").convert("L")
            candidates.append(bw)

        for img in candidates:
            for psm in (6, 4, 3, 11):
                text, conf, _ = run_ocr(img, psm=psm)
                score = score_ocr_text(text, conf=conf)
                if score > best_score:
                    best_score = score
                    best_text = text

        # Final safety net: rotate 90/180/270 on downscaled
        if best_score < 18:
            for angle in (90, 180, 270):
                rotated = fast.rotate(angle, expand=True)
                text, conf, _ = run_ocr(rotated, psm=6)
                score = score_ocr_text(text, conf=conf)
                if score > best_score:
                    best_score = score
                    best_text = text

        if best_score < 22 and os.environ.get("OCR_REGION_MODE", "true").strip().lower() in {"1", "true", "yes"}:
            region_text = ocr_text_from_regions(gray_full, lang=lang)
            region_score = score_ocr_text(region_text)
            if region_score > best_score:
                best_text = region_text
                best_score = region_score

        return best_text, best_score

    base = image
    best_text, best_score = ocr_from_base(base)

    if os.environ.get("OCR_CARD_MODE", "true").strip().lower() in {"1", "true", "yes"}:
        try:
            card_view = extract_document_perspective(base)
        except Exception:
            card_view = None
        if card_view is not None:
            card_text, card_score = ocr_from_base(card_view)
            if card_score > best_score:
                best_text, best_score = card_text, card_score

    return best_text

def post_process_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", "")
    # Repair obvious email/URL spacing issues.
    def collapse_email(match):
        return re.sub(r"\s+", "", match.group(0))
    text = re.sub(
        r"[A-Za-z0-9._%+-][A-Za-z0-9._%+\-\s]*@\s*[A-Za-z0-9.\-\s]+\.[A-Za-z]{2,}",
        collapse_email,
        text,
    )
    text = re.sub(
        r"https?://[A-Za-z0-9./%_+\-\s]+",
        lambda m: m.group(0).replace(" ", ""),
        text,
    )
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()

def extract_text_from_office_media(path: Path, media_prefix: str):
    text = ""
    if not path.is_file():
        return text
    try:
        with zipfile.ZipFile(path, "r") as archive:
            for info in archive.infolist():
                if not info.filename.lower().startswith(media_prefix):
                    continue
                if not info.filename.lower().endswith((".png", ".jpg", ".jpeg", ".tif", ".tiff")):
                    continue
                with archive.open(info) as image_file:
                    image = Image.open(image_file)
                    text += ocr_image_to_text(image)
    except zipfile.BadZipFile:
        return text
    return text

def extract_text_from_pdf(path):
    force_ocr = os.environ.get("FORCE_PDF_OCR", "").strip().lower() in {"1", "true", "yes"}
    text = ""
    if not force_ocr:
        with warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always", category=PyPDF2.errors.PdfReadWarning)
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        seen_messages = set()
        for warning in caught:
            message = str(warning.message)
            if message in seen_messages:
                continue
            seen_messages.add(message)
            print(f"PdfReadWarning [{Path(path)}]: {message}")

    def is_low_quality_text(value: str) -> bool:
        sample = value[:8000]
        if len(sample) < 200:
            return False
        if "\x00" in sample:
            return True
        tokens = re.findall(r"\S+", sample)
        if not tokens:
            return False
        weird_token_count = 0
        for token in tokens:
            if re.search(r"[\\^`\\[\\]{}]", token):
                weird_token_count += 1
                continue
            if re.search(r"[^A-Za-z0-9.,;:'\"()\\/&%$@#!?\-]", token):
                weird_token_count += 1
        if (weird_token_count / len(tokens)) > 0.15:
            return True
        weird_chars = sum(1 for ch in sample if ch in "\\[]{}^`")
        if (weird_chars / len(sample)) > 0.01:
            return True

        alpha_tokens = [t for t in tokens if t.isalpha()]
        if alpha_tokens:
            single_alpha_ratio = sum(1 for t in alpha_tokens if len(t) == 1) / len(alpha_tokens)
            if single_alpha_ratio > 0.08:
                return True
            split_pairs = 0
            for idx in range(len(tokens) - 1):
                left = tokens[idx]
                right = tokens[idx + 1]
                if len(left) == 1 and left.isalpha() and right.isalpha() and right[0].islower():
                    split_pairs += 1
            if split_pairs >= 3:
                return True

        if re.search(r"(?:\b[A-Za-z]\s+){3,}[A-Za-z]\b", sample):
            return True
        return False

    if text.strip() and not is_low_quality_text(text):
        return post_process_text(text)

    if text.strip() and not force_ocr:
        print(f"PDF OCR fallback [{Path(path)}]: low_quality_text")
        text = ""

    images = convert_from_path(path, dpi=300)
    for img in images:
        text += ocr_image_to_text(img)

    return post_process_text(text)


def extract_text_from_docx(path):
    doc = docx.Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    text += extract_text_from_office_media(Path(path), "word/media/")
    return text


def extract_text_from_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += f"{cell} "
    return text


def extract_text_from_image(path):
    img = Image.open(path)
    img = select_best_frame(img)
    exif_text = extract_exif_text(img)
    ocr_text = ocr_image_to_text(img, source_path=Path(path))
    combined = f"{exif_text}\n{ocr_text}" if exif_text else ocr_text
    return post_process_text(combined)

def extract_text_from_pptx(path):
    presentation = Presentation(path)
    slides_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    text = "\n".join(slides_text)
    text += extract_text_from_office_media(Path(path), "ppt/media/")
    return text

def strip_html(text):
    if not text:
        return ""
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
    text = re.sub(r"(?is)<.*?>", " ", text)
    return html_lib.unescape(text)

def extract_text_from_attachment_bytes(filename, content, depth):
    if depth > MAX_NESTED_DEPTH:
        return ""
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return ""
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(content)
        tmp.flush()
        return extract_text_from_file_path(Path(tmp.name), depth=depth)

def extract_text_from_eml_bytes(content, depth):
    if depth > MAX_NESTED_DEPTH:
        return ""
    msg = BytesParser(policy=policy.default).parsebytes(content)
    parts = []
    for part in msg.walk():
        if part.is_multipart():
            continue
        content_disposition = part.get_content_disposition()
        filename = part.get_filename()
        content_type = part.get_content_type()

        if content_disposition in {"attachment", "inline"} or filename:
            if filename:
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(
                        extract_text_from_attachment_bytes(
                            filename, payload, depth + 1
                        )
                    )
            continue

        if content_type == "text/plain":
            try:
                parts.append(part.get_content())
            except (LookupError, UnicodeDecodeError):
                payload = part.get_payload(decode=True) or b""
                parts.append(payload.decode(errors="ignore"))
        elif content_type == "text/html":
            try:
                parts.append(strip_html(part.get_content()))
            except (LookupError, UnicodeDecodeError):
                payload = part.get_payload(decode=True) or b""
                parts.append(strip_html(payload.decode(errors="ignore")))

    return "\n".join(p for p in parts if p)

def extract_text_from_eml_path(path, depth):
    with open(path, "rb") as f:
        return extract_text_from_eml_bytes(f.read(), depth=depth)

def extract_text_from_msg_path(path, depth):
    if depth > MAX_NESTED_DEPTH:
        return ""
    msg = extract_msg.Message(path)
    parts = []
    try:
        if msg.body:
            parts.append(msg.body)
        elif msg.htmlBody:
            parts.append(strip_html(msg.htmlBody))
        with tempfile.TemporaryDirectory() as tmpdir:
            for attachment in msg.attachments:
                filename = (
                    attachment.longFilename
                    or attachment.shortFilename
                    or attachment.filename
                )
                if not filename:
                    continue
                saved_path = attachment.save(customPath=tmpdir)
                if isinstance(saved_path, (tuple, list)):
                    saved_path = next(
                        (item for item in saved_path if isinstance(item, (str, bytes))),
                        None,
                    )
                if saved_path:
                    parts.append(
                        extract_text_from_file_path(
                            Path(saved_path), depth=depth + 1
                        )
                    )
    finally:
        msg.close()

    return "\n".join(p for p in parts if p)


def extract_text_from_file_path(path: Path, depth=0):
    ext = path.suffix.lower()
    text = ""
    if ext == ".pdf":
        text = extract_text_from_pdf(str(path))
    elif ext == ".docx":
        text = extract_text_from_docx(str(path))
    elif ext == ".xlsx":
        text = extract_text_from_xlsx(str(path))
    elif ext == ".pptx":
        text = extract_text_from_pptx(str(path))
    elif ext in {".jpg", ".jpeg", ".png", ".tif", ".tiff"}:
        text = extract_text_from_image(str(path))
    elif ext == ".eml":
        text = extract_text_from_eml_path(str(path), depth=depth)
    elif ext == ".msg":
        text = extract_text_from_msg_path(str(path), depth=depth)
    return post_process_text(text)


def extract_text_from_upload(file):
    suffix = Path(file.filename).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        file.save(tmp.name)
        return extract_text_from_file_path(Path(tmp.name))

def iter_documents_from_path(path: Path, depth=0, container_path=None, display_name=None):
    if depth > MAX_NESTED_DEPTH:
        return
    ext = path.suffix.lower()
    archive_exts = {".zip", ".7z", ".rar", ".tar", ".tgz", ".tar.gz"}
    if ext in archive_exts:
        container_label = container_path or display_name or str(path)
        yield from iter_documents_from_archive(
            path,
            container_label=container_label,
            depth=depth,
        )
        return
    text = extract_text_from_file_path(path, depth=depth)
    yield {
        "document": display_name or str(path),
        "text": text,
        "container_path": container_path,
        "inner_path": None,
        "source_path": path,
        "temp_path": False,
    }

def iter_documents_from_archive(path: Path, container_label: str, depth=0):
    ext = path.suffix.lower()
    if ext == ".zip":
        with zipfile.ZipFile(path, "r") as archive:
            for info in archive.infolist():
                if info.is_dir():
                    continue
                inner_name = info.filename
                with archive.open(info) as inner_file:
                    yield from iter_documents_from_archive_bytes(
                        inner_name,
                        inner_file.read(),
                        container_label,
                        depth + 1,
                    )
        return
    if ext in {".tar", ".tgz", ".tar.gz"}:
        mode = "r:gz" if ext in {".tgz", ".tar.gz"} else "r"
        with tarfile.open(path, mode) as archive:
            for member in archive.getmembers():
                if not member.isfile():
                    continue
                inner_file = archive.extractfile(member)
                if inner_file is None:
                    continue
                inner_name = member.name
                yield from iter_documents_from_archive_bytes(
                    inner_name,
                    inner_file.read(),
                    container_label,
                    depth + 1,
                )
        return
    if ext == ".rar":
        with rarfile.RarFile(path) as archive:
            for info in archive.infolist():
                if info.isdir():
                    continue
                inner_name = info.filename
                with archive.open(info) as inner_file:
                    yield from iter_documents_from_archive_bytes(
                        inner_name,
                        inner_file.read(),
                        container_label,
                        depth + 1,
                    )
        return
    if ext == ".7z":
        with py7zr.SevenZipFile(path, "r") as archive:
            for inner_name, bio in archive.readall().items():
                data = bio.read()
                yield from iter_documents_from_archive_bytes(
                    inner_name,
                    data,
                    container_label,
                    depth + 1,
                )
        return

def iter_documents_from_archive_bytes(inner_name, data, container_label, depth):
    if depth > MAX_NESTED_DEPTH:
        return
    suffix = Path(inner_name).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        return
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp.flush()
        yield {
            "document": f"{container_label}:{inner_name}",
            "text": extract_text_from_file_path(Path(tmp.name), depth=depth),
            "container_path": container_label,
            "inner_path": inner_name,
            "source_path": Path(tmp.name),
            "temp_path": True,
        }

# ======================================================
# Normalization helpers
# ======================================================

def normalize_match_text(value):
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip()

# ======================================================
# Proximity logic
# ======================================================

def nearest_regex_info(start, end, regex_hits):
    closest = None
    min_distance = None

    for r in regex_hits:
        if r["end"] < start:
            d = start - r["end"]
        elif r["start"] > end:
            d = r["start"] - end
        else:
            d = 0

        if min_distance is None or d < min_distance:
            min_distance = d
            closest = r["name"]

    return closest, min_distance

def build_rule_pack_xml(sit_outputs, sit_config, matched_patterns, rule_pack_bytes=None):
    regex_xml_entries = []
    for r in matched_patterns:
        regex_name = html_lib.escape(r.get("name", "Regex"))
        regex_pattern = html_lib.escape(r.get("pattern", ""))
        regex_xml_entries.append(
            f'    <Regex name="{regex_name}" pattern="{regex_pattern}" />'
        )
    regex_xml = "\n".join(regex_xml_entries) if regex_xml_entries else "    <!-- No regex patterns matched -->"

    if rule_pack_bytes:
        root = ET.fromstring(rule_pack_bytes)
        rule_pack_info = root.find("RulePackInfo")
        if rule_pack_info is None:
            rule_pack_info = ET.SubElement(root, "RulePackInfo")
        if sit_outputs["rule_pack_name"]:
            rule_pack_info.set("name", sit_outputs["rule_pack_name"])
        if sit_outputs["sit_publisher"]:
            rule_pack_info.set("publisher", sit_outputs["sit_publisher"])
        if sit_config.get("increment_rule_pack_version"):
            current_version = rule_pack_info.get("version", "")
            parts = current_version.split(".") if current_version else []
            if parts and all(part.isdigit() for part in parts):
                parts[-1] = str(int(parts[-1]) + 1)
                rule_pack_info.set("version", ".".join(parts))

        sit_parent = root.find("SensitiveInformationTypes")
        if sit_parent is None:
            sit_parent = ET.SubElement(root, "SensitiveInformationTypes")

        existing = None
        for child in sit_parent.findall("SensitiveInformationType"):
            if child.get("name") == sit_outputs["sit_name"]:
                existing = child
                break

        if sit_config.get("rule_pack_mode") == "update":
            if existing is None:
                raise ValueError(f"SIT not found in rule pack: {sit_outputs['sit_name']}")
            target = existing
        else:
            if existing is not None:
                raise ValueError(f"SIT already exists in rule pack: {sit_outputs['sit_name']}")
            target = ET.SubElement(
                sit_parent, "SensitiveInformationType", {"name": sit_outputs["sit_name"]}
            )

        if sit_outputs["sit_description"]:
            description = target.find("Description")
            if description is None:
                description = ET.SubElement(target, "Description")
            description.text = sit_outputs["sit_description"]

        regexes = target.find("Regexes")
        if regexes is None:
            regexes = ET.SubElement(target, "Regexes")
        else:
            regexes.clear()
        if matched_patterns:
            for r in matched_patterns:
                ET.SubElement(
                    regexes,
                    "Regex",
                    {
                        "name": r.get("name", "Regex"),
                        "pattern": r.get("pattern", ""),
                    },
                )
        else:
            regexes.append(ET.Comment("No regex patterns matched"))

        xml_text = ET.tostring(
            root, encoding="utf-8", xml_declaration=True
        ).decode("utf-8")
        version = rule_pack_info.get("version", "1.0")
        return xml_text, version

    xml_text = (
        "<?xml version=\"1.0\" encoding=\"utf-8\"?>\n"
        "<RulePack>\n"
        f"  <RulePackInfo name=\"{html_lib.escape(sit_outputs['rule_pack_name'])}\" "
        f"publisher=\"{html_lib.escape(sit_outputs['sit_publisher'])}\" version=\"1.0\" />\n"
        "  <SensitiveInformationTypes>\n"
        f"    <SensitiveInformationType name=\"{html_lib.escape(sit_outputs['sit_name'])}\">\n"
        f"      <Description>{html_lib.escape(sit_outputs['sit_description'])}</Description>\n"
        "      <Regexes>\n"
        f"{regex_xml}\n"
        "      </Regexes>\n"
        "    </SensitiveInformationType>\n"
        "  </SensitiveInformationTypes>\n"
        "</RulePack>\n"
    )
    return xml_text, "1.0"

# ======================================================
# API endpoint
# ======================================================

@app.route("/scan", methods=["POST"])
def scan():
    lexicon_path = request.form.get("lexicon_path", DEFAULT_LEXICON_PATH)
    regex_path = request.form.get("regex_path", REGEX_PATTERNS_PATH)
    scenario_mode = request.form.get("scenario_mode", "sit").strip().lower()
    ocr_backend = request.form.get("ocr_backend")
    lexicon_types = parse_lexicon_types(request.form.get("lexicon_types", ""))
    scan_policy_id = request.form.get("scan_policy_id")
    scan_id = request.form.get("scan_id")
    scan_path = request.form.get("path")
    recursive = request.form.get("recursive", "true").lower() == "true"
    log_to_stdout = request.form.get("log_stdout", "true").lower() == "true"
    debug_text = request.form.get("debug_text", "false").lower() == "true"
    debug_text_lines = []
    log_lines = []
    allowed_exts = parse_file_types(request.form.get("file_types"))
    output_path = request.form.get("output_path")
    batch_size = int(request.form.get("batch_size", "500"))
    return_limit = int(request.form.get("return_limit", "0"))

    if not lexicon_path:
        lexicon_path = DEFAULT_LEXICON_PATH
    if not regex_path:
        regex_path = REGEX_PATTERNS_PATH

    if scenario_mode not in {"sit", "scan_only"}:
        return jsonify({
            "success": False,
            "error": "scenario_mode must be 'sit' or 'scan_only'"
        }), 400

    if ocr_backend:
        os.environ["OCR_BACKEND"] = ocr_backend.strip().lower()
    elif scenario_mode == "scan_only" and os.environ.get("OCR_BACKEND", "auto").strip().lower() == "auto":
        # Favor speed for large batch scans unless explicitly overridden.
        os.environ["OCR_BACKEND"] = "tesseract"

    conn = get_db_connection()
    scan_params = {
        "lexicon_path": lexicon_path,
        "regex_path": regex_path,
        "scan_policy_id": scan_policy_id,
        "scan_id": scan_id,
        "scan_path": scan_path,
        "recursive": recursive,
        "log_to_stdout": log_to_stdout,
        "debug_text": debug_text,
        "allowed_exts": sorted(list(allowed_exts)) if allowed_exts else None,
        "output_path": output_path,
        "batch_size": batch_size,
        "return_limit": return_limit,
        "ocr_backend": ocr_backend or os.environ.get("OCR_BACKEND", "auto"),
    }
    cursor = conn.cursor()
    cursor.execute(
        """
        INSERT INTO scans (created_at, scenario_mode, status, scan_params)
        VALUES (?, ?, ?, ?)
        """,
        (
            datetime.now(timezone.utc).isoformat(),
            scenario_mode,
            "running",
            json.dumps(scan_params),
        ),
    )
    db_scan_id = cursor.lastrowid
    conn.commit()

    def set_scan_status(status):
        cursor.execute(
            "UPDATE scans SET status = ? WHERE id = ?",
            (status, db_scan_id),
        )
        conn.commit()

    progress_path = None
    if scan_id and re.fullmatch(r"[A-Za-z0-9_-]{6,64}", scan_id):
        PROGRESS_DIR.mkdir(parents=True, exist_ok=True)
        progress_path = PROGRESS_DIR / f"{scan_id}.json"

    def write_progress(payload):
        if not progress_path:
            return
        payload["scan_id"] = scan_id
        with open(progress_path, "w", encoding="utf-8") as f:
            json.dump(payload, f)

    if progress_path:
        write_progress({
            "status": "starting",
            "current": 0,
            "total": 0,
            "percent": 0,
            "current_document": "Listing files..."
        })

    temp_files = []

    def cleanup_temp_files():
        for path in temp_files:
            try:
                os.remove(path)
            except OSError:
                pass

    def save_uploaded_file(upload, suffix):
        if not upload or not upload.filename:
            return None
        ext = Path(upload.filename).suffix.lower()
        final_suffix = ext if ext else suffix
        with tempfile.NamedTemporaryFile(delete=False, suffix=final_suffix) as tmp:
            upload.save(tmp.name)
            temp_files.append(tmp.name)
            return tmp.name

    uploaded_lexicon = request.files.get("lexicon_file")
    uploaded_regex = request.files.get("regex_file")
    lexicon_upload_path = save_uploaded_file(uploaded_lexicon, ".csv")
    regex_upload_path = save_uploaded_file(uploaded_regex, ".py")
    if lexicon_upload_path:
        lexicon_path = lexicon_upload_path
    if regex_upload_path:
        regex_path = regex_upload_path

    scan_items = []
    listing_count = 0
    for scan_item in iter_scan_items(
        request.files.getlist("files"),
        scan_path,
        recursive,
        allowed_exts,
    ):
        scan_items.append(scan_item)
        listing_count += 1
        if scan_path and listing_count % 50 == 0:
            write_progress({
                "status": "listing",
                "current": listing_count,
                "total": 0,
                "percent": 0,
                "current_document": f"Found {listing_count} files..."
            })
    total_documents = listing_count
    if scenario_mode == "sit" and (total_documents == 0 or total_documents > MAX_SIT_FILES):
        write_progress({
            "status": "error",
            "error": f"Scenario 1 requires 1-{MAX_SIT_FILES} files. Found {total_documents}."
        })
        cleanup_temp_files()
        set_scan_status("error")
        conn.close()
        return jsonify({
            "success": False,
            "error": f"Scenario 1 requires 1-{MAX_SIT_FILES} files. Found {total_documents}.",
            "scan_log": "\n".join(log_lines),
            "debug_text": "\n".join(debug_text_lines) if debug_text else "",
        }), 400
    write_progress({
        "status": "running" if total_documents > 0 else "complete",
        "current": 0,
        "total": total_documents,
        "percent": 0 if total_documents > 0 else 100,
        "current_document": "" if total_documents > 0 else "No files found"
    })

    def log(line: str):
        if log_to_stdout:
            print(line)
        log_lines.append(line)

    def log_debug_text(line: str):
        if debug_text:
            debug_text_lines.append(line)

    if not os.path.isfile(regex_path):
        write_progress({
            "status": "error",
            "error": f"regex_path not found: {regex_path}"
        })
        cleanup_temp_files()
        set_scan_status("error")
        conn.close()
        return jsonify({
            "success": False,
            "error": f"regex_path not found: {regex_path}",
            "scan_log": "\n".join(log_lines),
            "debug_text": "\n".join(debug_text_lines) if debug_text else "",
        }), 400

    keywords = []
    keyword_set = set()
    if scenario_mode == "sit":
        keywords = load_keywords(lexicon_path, lexicon_types)
        keyword_set = set(keywords)
    regex_patterns = load_regex_patterns(regex_path)
    upsert_regex_patterns(conn, regex_patterns, regex_path)
    if scenario_mode == "sit":
        upsert_lexicon_entries(conn, load_lexicon_entries(lexicon_path), lexicon_path)
    regex_pattern_map = {p.get("name", "Regex"): p.get("pattern", "") for p in regex_patterns}
    # Keep debug_text focused on extracted document text only.

    sit_outputs = None
    sit_config = None
    if scenario_mode == "sit":
        def increment_version(value: str) -> str:
            if not value:
                return "1.0"
            parts = value.split(".")
            if not all(part.isdigit() for part in parts):
                return value
            parts[-1] = str(int(parts[-1]) + 1)
            return ".".join(parts)

        sit_outputs = {
            "sit_name": request.form.get("sit_name", "").strip() or "Custom SIT",
            "sit_description": request.form.get("sit_description", "").strip(),
            "sit_publisher": request.form.get("sit_publisher", "").strip() or "Purview Custom",
            "rule_pack_name": request.form.get("rule_pack_name", "").strip() or "CustomRulePack",
        }
        sit_config = {
            "rule_pack_mode": request.form.get("rule_pack_mode", "new").strip().lower(),
            "increment_rule_pack_version": (
                request.form.get("increment_rule_pack_version", "true").strip().lower()
                == "true"
            ),
            "rule_pack_base64": request.form.get("rule_pack_base64", "").strip(),
            "rule_pack_file": request.files.get("rule_pack_file"),
        }

    output_file = None
    if output_path:
        output_file = open(output_path, "a", encoding="utf-8")

    results = []
    batch_results = []
    errors = []
    documents_scanned = 0
    total_matches = 0
    regex_summary = {}
    keyword_summary = {}

    def record_match(item):
        nonlocal total_matches
        total_matches += 1

        if output_file:
            batch_results.append(item)
            if len(batch_results) >= batch_size:
                output_file.write("\n".join(json.dumps(r) for r in batch_results))
                output_file.write("\n")
                output_file.flush()
                batch_results.clear()

        if not output_file or (return_limit > 0 and len(results) < return_limit):
            results.append(item)

    for source, item in scan_items:
        documents_scanned += 1
        try:
            if source == "upload":
                temp_path = save_uploaded_file(item, Path(item.filename).suffix.lower() or ".bin")
                if not temp_path:
                    continue
                doc_entries = iter_documents_from_path(
                    Path(temp_path),
                    display_name=item.filename,
                )
                source_type = "upload"
                source_uri = item.filename
            else:
                doc_entries = iter_documents_from_path(item)
                source_type = "path"
                source_uri = str(item)

            for doc_entry in doc_entries:
                if doc_entry["temp_path"]:
                    temp_files.append(str(doc_entry["source_path"]))
                document = doc_entry["document"]
                text = doc_entry["text"]
                container_path = doc_entry["container_path"]
                inner_path = doc_entry["inner_path"]
                source_path = doc_entry["source_path"]

                if total_documents > 0:
                    write_progress({
                        "status": "running",
                        "current": documents_scanned,
                        "total": total_documents,
                        "percent": int((documents_scanned / total_documents) * 100),
                        "current_document": document
                    })

                if debug_text:
                    log_debug_text(f"started processing file: {document}")
                    log_debug_text(f"extracted_text_begin: {document}")
                    log_debug_text(text)
                    log_debug_text(f"extracted_text_end: {document}")
                    log_debug_text(f"ended processing file: {document}")

                file_sha256 = sha256_for_path(source_path)
                cursor.execute(
                    """
                    INSERT INTO documents
                        (scan_id, source_type, source_uri, sha256, extracted_text_ref)
                    VALUES (?, ?, ?, ?, ?)
                    """,
                    (
                        db_scan_id,
                        source_type,
                        source_uri if inner_path is None else document,
                        file_sha256,
                        None,
                    ),
                )
                document_id = cursor.lastrowid
                conn.commit()

                # -------------------------
                # Regex hits
                # -------------------------
                regex_hits = []
                for r in regex_patterns:
                    for m in re.finditer(r["pattern"], text, re.IGNORECASE | re.DOTALL):
                        regex_hits.append({
                            "name": r["name"],
                            "start": m.start(),
                            "end": m.end(),
                            "value": m.group(0)
                        })

                # -------------------------
                # Keyword hits (strict), only after a regex match exists
                # -------------------------
                keyword_hits = []
                if regex_hits and scenario_mode == "sit":
                    token_iter = list(re.finditer(r"[A-Za-z0-9]+", text))
                    for idx in range(len(token_iter) - 1):
                        first = token_iter[idx]
                        second = token_iter[idx + 1]
                        between = text[first.end():second.start()]
                        if between != " ":
                            continue
                        w1 = first.group(0).lower()
                        w2 = second.group(0).lower()
                        if w1 in keyword_set and w2 in keyword_set:
                            keyword_hits.append({
                                "phrase": f"{w1} {w2}",
                                "start": first.start(),
                                "end": second.end()
                            })

                logged_file = False
                for r in regex_hits:
                    if not logged_file:
                        log(f"matched_file: {document}")
                        logged_file = True
                    log(
                        f"match: regex={r['name']} "
                        f"document={document} "
                        f"value={r['value']}"
                    )
                    match_payload = {
                        "type": "regex",
                        "fallback": False,
                        "regex_name": r["name"],
                        "document": document,
                        "position": r["start"],
                        "value": r["value"],
                        "context": text[max(0, r["start"]-50):r["end"]+50],
                        "container_path": container_path,
                        "inner_path": inner_path,
                    }
                    record_match(match_payload)
                    cursor.execute(
                        """
                        INSERT INTO matches (
                            scan_id, document_id, match_type, pattern_name, value,
                            position, context, nearest_regex, distance, container_path, inner_path
                        )
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (
                            db_scan_id,
                            document_id,
                            "regex",
                            r["name"],
                            r["value"],
                            r["start"],
                            match_payload["context"],
                            None,
                            None,
                            container_path,
                            inner_path,
                        ),
                    )
                    conn.commit()
                    if scenario_mode == "sit":
                        key = r["name"]
                        summary = regex_summary.setdefault(
                            key,
                            {
                                "regex_name": r["name"],
                                "pattern": regex_pattern_map.get(r["name"], ""),
                                "keyword_match_count": 0,
                                "total_count": 0,
                                "documents": set(),
                                "matched_texts": set(),
                            },
                        )
                        summary["total_count"] += 1
                        summary["documents"].add(document)
                        normalized_value = normalize_match_text(r["value"])
                        if normalized_value:
                            summary["matched_texts"].add(normalized_value)

                for k in keyword_hits:
                    if not logged_file:
                        log(f"matched_file: {document}")
                        logged_file = True
                    nearest_name = None
                    nearest_dist = None
                    if regex_hits:
                        nearest_name, nearest_dist = nearest_regex_info(
                            k["start"], k["end"], regex_hits
                        )
                    log(
                        f"match: keyword={k['phrase']} "
                        f"document={document}"
                    )
                    match_payload = {
                        "type": "keyword",
                        "fallback": False,
                        "phrase": k["phrase"],
                        "document": document,
                        "position": k["start"],
                        "nearest_regex": nearest_name,
                        "nearest_regex_distance": nearest_dist,
                        "context": text[max(0, k["start"]-50):k["end"]+50],
                        "container_path": container_path,
                        "inner_path": inner_path,
                    }
                    record_match(match_payload)
                    cursor.execute(
                        """
                        INSERT INTO matches (
                            scan_id, document_id, match_type, pattern_name, value,
                            position, context, nearest_regex, distance, container_path, inner_path
                        )
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (
                            db_scan_id,
                            document_id,
                            "keyword",
                            k["phrase"],
                            None,
                            k["start"],
                            match_payload["context"],
                            nearest_name,
                            nearest_dist,
                            container_path,
                            inner_path,
                        ),
                    )
                    conn.commit()
                    if scenario_mode == "sit":
                        key = k["phrase"]
                        summary = keyword_summary.setdefault(
                            key,
                            {
                                "phrase": k["phrase"],
                                "total_count": 0,
                                "documents": set(),
                                "triggering_regexes": set(),
                            },
                        )
                        summary["total_count"] += 1
                        summary["documents"].add(document)
                        if nearest_name:
                            summary["triggering_regexes"].add(nearest_name)
                            regex_summary.setdefault(
                                nearest_name,
                                {
                                    "regex_name": nearest_name,
                                    "pattern": regex_pattern_map.get(nearest_name, ""),
                                    "keyword_match_count": 0,
                                    "total_count": 0,
                                    "documents": set(),
                                    "matched_texts": set(),
                                },
                            )["keyword_match_count"] += 1

        except Exception as e:
            error_doc = document if "document" in locals() else "unknown"
            errors.append(f"{error_doc}: {str(e)}")
            log(f"error: document={error_doc} error={e}")

    if output_file and batch_results:
        output_file.write("\n".join(json.dumps(r) for r in batch_results))
        output_file.write("\n")
        output_file.flush()
        batch_results.clear()

    if output_file:
        output_file.close()

    regex_summary_list = []
    keyword_summary_list = []
    if scenario_mode == "sit":
        for summary in regex_summary.values():
            docs = sorted(summary["documents"])
            regex_summary_list.append({
                "regex_name": summary["regex_name"],
                "pattern": summary["pattern"],
                "multi_phrase_keyword_count": summary.get("keyword_match_count", 0),
                "total_count": summary["total_count"],
                "file_count": len(docs),
                "documents": docs,
                "matched_text": sorted(summary.get("matched_texts", set())),
                "priority": (len(docs) * 1000) + summary["total_count"],
            })
        for summary in keyword_summary.values():
            docs = sorted(summary["documents"])
            triggering = sorted(summary.get("triggering_regexes", set()))
            keyword_summary_list.append({
                "phrase": summary["phrase"],
                "total_count": summary["total_count"],
                "file_count": len(docs),
                "documents": docs,
                "triggering_regex": triggering,
                "priority": (len(docs) * 1000) + summary["total_count"],
            })

        regex_summary_list.sort(
            key=lambda item: (-item["file_count"], -item["total_count"], item["regex_name"])
        )
        keyword_summary_list.sort(
            key=lambda item: (-item["file_count"], -item["total_count"], item["phrase"])
        )

        if sit_outputs and sit_config is not None:
            matched_regex_names = {item["regex_name"] for item in regex_summary_list}
            matched_patterns = [
                r for r in regex_patterns if r.get("name") in matched_regex_names
            ]

            uploaded_rule_pack = sit_config["rule_pack_file"]
            rule_pack_base64 = sit_config["rule_pack_base64"]
            rule_pack_bytes = None
            if uploaded_rule_pack and uploaded_rule_pack.filename:
                rule_pack_bytes = uploaded_rule_pack.read()
            elif rule_pack_base64:
                try:
                    rule_pack_bytes = base64.b64decode(rule_pack_base64)
                except (ValueError, TypeError) as exc:
                    set_scan_status("error")
                    conn.close()
                    return jsonify({
                        "success": False,
                        "error": f"invalid rule_pack_base64: {exc}",
                        "scan_log": "\n".join(log_lines),
                        "debug_text": "\n".join(debug_text_lines) if debug_text else "",
                    }), 400

            if rule_pack_bytes:
                try:
                    sit_outputs["rule_pack_xml"], sit_outputs["rule_pack_version"] = build_rule_pack_xml(
                        sit_outputs,
                        sit_config,
                        matched_patterns,
                        rule_pack_bytes=rule_pack_bytes,
                    )
                except (ET.ParseError, ValueError) as exc:
                    write_progress({
                        "status": "error",
                        "error": f"invalid rule pack xml: {exc}"
                    })
                    cleanup_temp_files()
                    set_scan_status("error")
                    conn.close()
                    return jsonify({
                        "success": False,
                        "error": f"invalid rule pack xml: {exc}",
                        "scan_log": "\n".join(log_lines),
                        "debug_text": "\n".join(debug_text_lines) if debug_text else "",
                    }), 400
            else:
                sit_outputs["rule_pack_xml"], sit_outputs["rule_pack_version"] = build_rule_pack_xml(
                    sit_outputs,
                    sit_config,
                    matched_patterns,
                )

            sit_outputs["powershell_script"] = (
                "$rulePackPath = \"./"
                + sit_outputs["rule_pack_name"]
                + ".xml\"\n"
                "Write-Host \"Importing rule pack: $rulePackPath\"\n"
                "New-DlpSensitiveInformationTypeRulePackage -FileData "
                "(Get-Content -Path $rulePackPath -Encoding Byte -ReadCount 0)\n"
            )

            cursor.execute(
                """
                INSERT INTO sit_definitions
                    (name, description, publisher, created_from_scan_id, rule_pack_name, version)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    sit_outputs["sit_name"],
                    sit_outputs["sit_description"],
                    sit_outputs["sit_publisher"],
                    db_scan_id,
                    sit_outputs["rule_pack_name"],
                    sit_outputs.get("rule_pack_version", "1.0"),
                ),
            )
            sit_definition_id = cursor.lastrowid
            cursor.execute(
                """
                INSERT INTO rule_packs
                    (sit_definition_id, xml_blob, generated_at, generation_params)
                VALUES (?, ?, ?, ?)
                """,
                (
                    sit_definition_id,
                    sit_outputs["rule_pack_xml"],
                    datetime.now(timezone.utc).isoformat(),
                    json.dumps({
                        "rule_pack_mode": sit_config.get("rule_pack_mode"),
                        "increment_rule_pack_version": sit_config.get("increment_rule_pack_version"),
                        "rule_pack_base64": bool(sit_config.get("rule_pack_base64")),
                        "rule_pack_file": (
                            sit_config.get("rule_pack_file").filename
                            if sit_config.get("rule_pack_file") else None
                        ),
                    }),
                ),
            )
            conn.commit()

    cleanup_temp_files()
    write_progress({
        "status": "complete",
        "current": documents_scanned,
        "total": total_documents,
        "percent": 100,
        "current_document": ""
    })
    set_scan_status("complete")
    conn.close()

    return jsonify({
        "success": True,
        "db_scan_id": db_scan_id,
        "scenario_mode": scenario_mode,
        "sit_outputs": sit_outputs,
        "documents_scanned": documents_scanned,
        "total_matches": total_matches,
        "matches": results,
        "regex_summary": regex_summary_list,
        "keyword_summary": keyword_summary_list,
        "errors": errors,
        "scan_log": "\n".join(log_lines),
        "debug_text": "\n".join(debug_text_lines) if debug_text else ""
    })


@app.route("/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "ocr_available": True
    })


@app.route("/", methods=["GET"])
def index():
    return send_from_directory(".", "index.html")


@app.route("/index.html", methods=["GET"])
def index_html():
    return send_from_directory(".", "index.html")

@app.route("/regex-files", methods=["GET"])
def regex_files():
    root = Path(".").resolve()
    files = []
    try:
        for path in list(root.glob("regex_patterns*.py")) + list(root.glob("regex_patterns*.json")):
            if path.is_file():
                files.append(path.name)
                if len(files) >= 200:
                    break
    except OSError:
        pass
    files.sort(key=lambda name: (0 if name.endswith(".py") else 1, name))
    return jsonify({"success": True, "files": files})


@app.route("/lexicon-files", methods=["GET"])
def lexicon_files():
    root = Path(".").resolve()
    files = []
    try:
        for path in root.glob("lexicon*.csv"):
            if path.is_file():
                files.append(path.name)
                if len(files) >= 200:
                    break
    except OSError:
        pass
    files.sort()
    return jsonify({"success": True, "files": files})

@app.route("/rulepack-cache", methods=["GET"])
def rulepack_cache():
    if not RULEPACK_CACHE_PATH.is_file():
        return jsonify({
            "success": False,
            "error": "rule pack cache not found"
        }), 404
    try:
        with open(RULEPACK_CACHE_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
    except (OSError, json.JSONDecodeError) as exc:
        return jsonify({
            "success": False,
            "error": f"unable to read rule pack cache: {exc}"
        }), 400
    return jsonify({"success": True, "cache": data})

@app.route("/rule-pack/export", methods=["POST"])
def export_rule_pack():
    payload = request.get_json(silent=True) or request.form
    sit_definition_id = payload.get("sit_definition_id")
    scan_id_value = payload.get("scan_id")

    conn = get_db_connection()
    cursor = conn.cursor()
    sit_row = None

    if sit_definition_id:
        cursor.execute(
            "SELECT * FROM sit_definitions WHERE id = ?",
            (sit_definition_id,),
        )
        sit_row = cursor.fetchone()
    elif scan_id_value:
        try:
            scan_id_int = int(scan_id_value)
        except (TypeError, ValueError):
            conn.close()
            return jsonify({
                "success": False,
                "error": "scan_id must be an integer database scan id",
            }), 400
        cursor.execute(
            """
            SELECT * FROM sit_definitions
            WHERE created_from_scan_id = ?
            ORDER BY id DESC
            LIMIT 1
            """,
            (scan_id_int,),
        )
        sit_row = cursor.fetchone()

    if sit_row is None:
        conn.close()
        return jsonify({
            "success": False,
            "error": "sit_definition_id or scan_id not found",
        }), 404

    cursor.execute(
        """
        SELECT * FROM rule_packs
        WHERE sit_definition_id = ?
        ORDER BY id DESC
        LIMIT 1
        """,
        (sit_row["id"],),
    )
    rule_pack_row = cursor.fetchone()
    if rule_pack_row:
        conn.close()
        return jsonify({
            "success": True,
            "sit_definition_id": sit_row["id"],
            "rule_pack_xml": rule_pack_row["xml_blob"],
        })

    cursor.execute(
        """
        SELECT DISTINCT pattern_name FROM matches
        WHERE scan_id = ? AND match_type = 'regex'
        """,
        (sit_row["created_from_scan_id"],),
    )
    matched_names = {row["pattern_name"] for row in cursor.fetchall()}
    matched_patterns = []
    if matched_names:
        cursor.execute(
            """
            SELECT name, pattern FROM regex_patterns
            WHERE name IN ({})
            """.format(",".join("?" for _ in matched_names)),
            tuple(matched_names),
        )
        matched_patterns = [dict(row) for row in cursor.fetchall()]

    sit_outputs = {
        "sit_name": sit_row["name"],
        "sit_description": sit_row["description"] or "",
        "sit_publisher": sit_row["publisher"] or "Purview Custom",
        "rule_pack_name": sit_row["rule_pack_name"],
    }
    sit_config = {
        "rule_pack_mode": "new",
        "increment_rule_pack_version": True,
    }
    rule_pack_xml, version = build_rule_pack_xml(
        sit_outputs,
        sit_config,
        matched_patterns,
    )
    cursor.execute(
        """
        INSERT INTO rule_packs
            (sit_definition_id, xml_blob, generated_at, generation_params)
        VALUES (?, ?, ?, ?)
        """,
        (
            sit_row["id"],
            rule_pack_xml,
            datetime.now(timezone.utc).isoformat(),
            json.dumps({"generated_by": "rule-pack-export"}),
        ),
    )
    conn.commit()
    conn.close()
    return jsonify({
        "success": True,
        "sit_definition_id": sit_row["id"],
        "rule_pack_xml": rule_pack_xml,
        "version": version,
    })

@app.route("/progress/<scan_id>", methods=["GET"])
def progress(scan_id):
    if not re.fullmatch(r"[A-Za-z0-9_-]{6,64}", scan_id):
        return jsonify({"success": False, "error": "invalid scan_id"}), 400
    progress_path = PROGRESS_DIR / f"{scan_id}.json"
    if not progress_path.is_file():
        return jsonify({"success": False, "error": "scan_id not found"}), 404
    try:
        with open(progress_path, "r", encoding="utf-8") as f:
            raw = f.read().strip()
        if not raw:
            data = {
                "status": "running",
                "percent": 0,
                "current": 0,
                "total": 0,
                "current_document": "",
            }
        else:
            data = json.loads(raw)
    except json.JSONDecodeError:
        data = {
            "status": "running",
            "percent": 0,
            "current": 0,
            "total": 0,
            "current_document": "",
        }
    data["success"] = True
    return jsonify(data)


if __name__ == "__main__":
    print("Starting Document Scanner Backend Server...")
    port = int(os.environ.get("PORT", "5000"))
    app.run(debug=True, use_reloader=False, port=port)
